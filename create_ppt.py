from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色方案 - 科技感蓝 + 活力橙
PRIMARY_BLUE = RgbColor(41, 128, 185)
DARK_BLUE = RgbColor(44, 62, 80)
ACCENT_ORANGE = RgbColor(230, 126, 34)
LIGHT_BG = RgbColor(236, 240, 241)
WHITE = RgbColor(255, 255, 255)

# ===== 第1页：封面 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 背景色块
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

# 主标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
tf = title_box.text_frame
tf.text = "元枢纽 Meta-Axis"
p = tf.paragraphs[0]
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 副标题
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1))
tf2 = subtitle_box.text_frame
tf2.text = "RC 攀爬车周末训练营"
p2 = tf2.paragraphs[0]
p2.font.size = Pt(40)
p2.font.color.rgb = ACCENT_ORANGE
p2.alignment = PP_ALIGN.CENTER

# 描述
desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(0.8))
tf3 = desc_box.text_frame
tf3.text = "教育部白名单赛事 · EF-Core 执行功能训练 · 从科普到竞技"
p3 = tf3.paragraphs[0]
p3.font.size = Pt(20)
p3.font.color.rgb = LIGHT_BG
p3.alignment = PP_ALIGN.CENTER

# ===== 第2页：痛点 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY_BLUE
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "您的孩子是否面临这些挑战？"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# 痛点列表
pain_points = [
    "做作业拖拖拉拉，总是最后一刻才开始",
    "容易分心，难以专注完成任务",
    "书包和书桌杂乱无章，经常找不到东西",
    "做决定不思考后果，容易冲动行事",
    "难以同时记住和处理多项信息",
    "不了解自己的学习风格，无法自我调整"
]

y_pos = 1.6
for point in pain_points:
    # 橙色圆点
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y_pos + 0.1), Inches(0.2), Inches(0.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_ORANGE
    circle.line.fill.background()
    
    # 文字
    text_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos), Inches(11), Inches(0.6))
    tf = text_box.text_frame
    tf.text = point
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_BLUE
    y_pos += 0.9

# 底部强调
emphasis_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
tf = emphasis_box.text_frame
tf.text = "这些问题的根源不是意志力薄弱，而是\"操作系统\"需要升级"
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE
p.alignment = PP_ALIGN.CENTER

# ===== 第3页：解决方案 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY_BLUE
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "元枢纽 EF-Core：重构心智操作系统"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# 核心理念
core_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(12), Inches(1))
tf = core_box.text_frame
tf.text = "大脑如电脑，智力是硬件，EF-Core 是操作系统(OS)"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# 六大维度
ef_dims = [
    ("任务启动", "克服拖延，立即开始"),
    ("时间体感", "校准主观时间与客观时间"),
    ("空间管理", "结构化整理，减少认知负荷"),
    ("优先级决策", "在干扰中识别核心价值"),
    ("持续专注", "构建深度学习的心流区"),
    ("元认知监控", "学会\"思考自己的思考\"")
]

y_pos = 2.8
for i, (title, desc) in enumerate(ef_dims):
    x_pos = 0.8 if i < 3 else 7
    y_offset = y_pos + (i % 3) * 1.5
    
    # 卡片背景
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(y_offset), Inches(5.5), Inches(1.3))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = PRIMARY_BLUE
    card.line.width = Pt(2)
    
    # 标题
    title_tb = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(y_offset + 0.1), Inches(5), Inches(0.5))
    tf = title_tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    
    # 描述
    desc_tb = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(y_offset + 0.6), Inches(5), Inches(0.6))
    tf = desc_tb.text_frame
    tf.text = desc
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BLUE

# ===== 第4页：项目权威性 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_ORANGE
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "教育部白名单赛事 · 权威认证"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# 赛事名称
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.8), Inches(9.333), Inches(1.2))
badge.fill.solid()
badge.fill.fore_color.rgb = PRIMARY_BLUE
badge.line.fill.background()

event_box = slide.shapes.add_textbox(Inches(2.5), Inches(2), Inches(8.5), Inches(1))
tf = event_box.text_frame
tf.text = '"驾驭未来"全国青少年车辆模型教育竞赛'
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 权威性要点
auth_points = [
    "教育部白名单在册体育运动项目",
    "完整的赛事组织委员会（组委会）",
    "小学/中学分组，男子/女子分组",
    "个人项目 + 团体项目双重赛道",
    "从科普入门到竞技精通的完整路径"
]

y_pos = 3.5
for point in auth_points:
    check = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(y_pos + 0.15), Inches(0.25), Inches(0.25))
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT_ORANGE
    check.line.fill.background()
    
    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(11), Inches(0.6))
    tf = text_box.text_frame
    tf.text = point
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_BLUE
    y_pos += 0.8

# ===== 第5页：10期课程总览 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY_BLUE
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "10期周末训练营：从科普到竞技"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# 三阶段说明
stages = [
    ("第1-3期", "系统重构期", "建立认知 · 规范操作 · 空间规划", LIGHT_BG),
    ("第4-7期", "技能内化期", "专注训练 · 团队协作 · 工程思维", RgbColor(174, 214, 241)),
    ("第8-10期", "独立运行期", "赛事备战 · 全维实战 · 冠军传承", RgbColor(133, 193, 233))
]

y_pos = 1.5
for stage_num, stage_name, stage_desc, color in stages:
    # 阶段卡片
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12.5), Inches(1.6))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.color.rgb = PRIMARY_BLUE
    card.line.width = Pt(2)
    
    # 期数
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.2), Inches(2), Inches(0.5))
    tf = num_box.text_frame
    tf.text = stage_num
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    
    # 阶段名
    name_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.7), Inches(3), Inches(0.5))
    tf = name_box.text_frame
    tf.text = stage_name
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # 描述
    desc_box = slide.shapes.add_textbox(Inches(4), Inches(y_pos + 0.5), Inches(8.5), Inches(0.8))
    tf = desc_box.text_frame
    tf.text = stage_desc
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_BLUE
    
    y_pos += 2

# ===== 第6页：核心能力培养 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_ORANGE
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "RC 车训练如何提升 EF 核心能力？"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# 能力对应表
abilities = [
    ("任务启动", "将复杂的车辆调试拆解为可执行的步骤", "克服拖延，立即开始"),
    ("时间体感", "预估圈时 vs 实际圈时的校准训练", "建立准确的时间观念"),
    ("空间管理", "赛道设计中的物理空间规划", "结构化思维，减少混乱"),
    ("优先级决策", "过障时的路径选择与风险权衡", "在压力下做出最优决策"),
    ("持续专注", "完整圈速挑战中的稳定发挥", "构建心流，深度投入"),
    ("元认知监控", "每次驾驶后的复盘与策略调整", "学会反思，持续改进")
]

y_pos = 1.5
for ability, rc_apply, result in abilities:
    # 能力名称（橙色标签）
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(2.2), Inches(0.6))
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT_ORANGE
    tag.line.fill.background()
    
    tag_text = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.1), Inches(2), Inches(0.5))
    tf = tag_text.text_frame
    tf.text = ability
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # RC应用
    apply_box = slide.shapes.add_textbox(Inches(3), Inches(y_pos), Inches(5.5), Inches(0.6))
    tf = apply_box.text_frame
    tf.text = rc_apply
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BLUE
    
    # 成果
    result_box = slide.shapes.add_textbox(Inches(8.8), Inches(y_pos), Inches(4), Inches(0.6))
    tf = result_box.text_frame
    tf.text = "→ " + result
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    
    y_pos += 0.9

# ===== 第7页：年度大项目 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 深色背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
tf = title_box.text_frame
tf.text = "年度成就项目"
p = tf.paragraphs[0]
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE
p.alignment = PP_ALIGN.CENTER

# 赛事名称
event_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(1))
tf = event_box.text_frame
tf.text = "元枢纽杯 RC 攀爬车锦标赛"
p = tf.paragraphs[0]
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 描述
desc_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.5), Inches(1))
tf = desc_box.text_frame
tf.text = "不是模拟考试，是真实世界的压力测试"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.color.rgb = LIGHT_BG
p.alignment = PP_ALIGN.CENTER

# 赛事要素
elements = [
    "资格赛：个人计时挑战",
    "正赛：多轮积分制竞技", 
    "团队赛：接力协作挑战",
    "障碍赛：越野通关挑战"
]

y_pos = 4.8
for elem in elements:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(y_pos + 0.15), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_ORANGE
    dot.line.fill.background()
    
    text_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(8), Inches(0.5))
    tf = text_box.text_frame
    tf.text = elem
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    y_pos += 0.7

# ===== 第8页：数据化成果 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY_BLUE
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "数据化成长追踪：看得见的能力提升"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# 雷达图说明（用文字描述）
radar_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.6))
tf = radar_title.text_frame
tf.text = "个人 EF 技能雷达图"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# 雷达图维度
dimensions = [
    "任务启动能力",
    "时间体感精度", 
    "空间管理效率",
    "优先级决策质量",
    "持续专注时长",
    "元认知监控水平"
]

y_pos = 2.2
for dim in dimensions:
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(5), Inches(0.5))
    tf = text_box.text_frame
    tf.text = "• " + dim
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    y_pos += 0.6

# 成长对比
compare_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(0.6))
tf = compare_title.text_frame
tf.text = "训练前后对比"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

comparisons = [
    ("任务启动", "30% → 85%"),
    ("时间预估准确率", "40% → 80%"),
    ("专注时长", "15min → 45min"),
    ("空间整理效率", "50% → 90%")
]

y_pos = 2.2
for skill, improvement in comparisons:
    # 技能名
    skill_box = slide.shapes.add_textbox(Inches(7), Inches(y_pos), Inches(3.5), Inches(0.5))
    tf = skill_box.text_frame
    tf.text = skill
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    
    # 提升数据
    data_box = slide.shapes.add_textbox(Inches(10.5), Inches(y_pos), Inches(2.5), Inches(0.5))
    tf = data_box.text_frame
    tf.text = improvement
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE
    
    y_pos += 0.7

# ===== 第9页：报名信息 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 标题背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_ORANGE
shape.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
tf = title_box.text_frame
tf.text = "报名信息"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

# 信息卡片
info_items = [
    ("课程周期", "10期周末训练营，每期1天（上午+下午各2小时）"),
    ("适合年龄", "小学组（6-12岁）· 中学组（13-15岁）"),
    ("班级规模", "小班教学，每班8-12人"),
    ("包含内容", "RC攀爬车设备使用 · 专业赛道 · 教练指导 · 赛事报名"),
    ("成果交付", "个人赛车档案 · EF技能雷达图 · 年度成就展示 · 锦标赛参赛资格")
]

y_pos = 1.5
for label, content in info_items:
    # 标签
    label_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(2.5), Inches(0.5))
    tf = label_box.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos), Inches(9), Inches(0.8))
    tf = content_box.text_frame
    tf.text = content
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    
    y_pos += 1.1

# 联系方式
contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.8))
tf = contact_box.text_frame
tf.text = "扫码咨询 · 免费体验课预约"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ACCENT_ORANGE
p.alignment = PP_ALIGN.CENTER

# ===== 第10页：封底 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 背景
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

# 品牌名
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
tf = title_box.text_frame
tf.text = "元枢纽 Meta-Axis"
p = tf.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Slogan
slogan_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(1))
tf = slogan_box.text_frame
tf.text = "重构心智中枢 · 掌舵未来航向"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = ACCENT_ORANGE
p.alignment = PP_ALIGN.CENTER

# 副标题
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(0.8))
tf = sub_box.text_frame
tf.text = "赋予孩子穿越未来迷雾的能力"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_BG
p.alignment = PP_ALIGN.CENTER

# 保存
output_path = '/Users/pengwang/.openclaw/workspace/元枢纽_RC课程招生PPT.pptx'
prs.save(output_path)
print(f"PPT 已保存：{output_path}")
print("\n共10页：")
print("1. 封面 - 品牌+主题")
print("2. 痛点 - 6大学习困难")
print("3. 解决方案 - EF-Core六大维度")
print("4. 权威性 - 教育部白名单赛事")
print("5. 课程总览 - 三阶段10期")
print("6. 能力培养 - RC训练如何提升EF")
print("7. 年度大项目 - 元枢纽杯锦标赛")
print("8. 数据化成果 - 雷达图+成长对比")
print("9. 报名信息 - 周期/年龄/内容")
print("10. 封底 - 品牌slogan")
